from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

def create_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[1] # 'Title and Content' layout
    slide = presentation.slides.add_slide(slide_layout)

    # Set font size
    slide.shapes.title.text = title
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.font.name = 'San Francisco'
    title_para.font.size = Pt(32)

    content_placeholder = slide.placeholders[1]

    tf = content_placeholder.text_frame
    tf.text = content[0]
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.name = 'San Francisco'
    for line in content[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = 'San Francisco'


def create_ppt(ppt_title, slide_contents, outfile_path="test_ppt.pptx"):
    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Set the title slide layout
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)

    # Add a title to the slide
    title = slide.shapes.title
    # title.font = Pt(14)
    title.text = ppt_title

    # Add content to the slides
    for slide_content in slide_contents:
        create_slide(presentation, slide_content['title'], slide_content['content'])

    # Save the PowerPoint
    presentation.save(outfile_path)
    # save presentation as binary output
    binary_output = BytesIO()
    presentation.save(binary_output)
    return binary_output



# slide_contents = [
#     {"title": "Slide 1", "content": ["Point 1", "Point 2", "Point 3"]},
#     {"title": "Slide 2", "content": ["Point 1", "Point 2", "Point 3"]},
# ]

# create_ppt("My Presentation", slide_contents)

